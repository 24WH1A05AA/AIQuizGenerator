"""PPT/PPTX text extraction using python-pptx."""

import hashlib
import io

import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


# ── Exceptions ─────────────────────────────────────────────────────────────

class PPTParseError(Exception):
    """Base class for all PPT parsing errors."""


class UnsupportedFormatError(PPTParseError):
    """File is not a supported PPTX format (e.g. legacy .ppt, PDF, DOCX)."""


class CorruptedFileError(PPTParseError):
    """File has the right magic bytes but cannot be opened or parsed."""


class EmptyPresentationError(PPTParseError):
    """Presentation opened successfully but contains no readable text."""


# ── Format detection ───────────────────────────────────────────────────────

_MAGIC_PPTX = b"PK\x03\x04"        # ZIP / Office Open XML (.pptx, .docx, …)
_MAGIC_PPT  = b"\xd0\xcf\x11\xe0"  # OLE Compound Document (legacy .ppt)


def _detect_format(data: bytes) -> str:
    """Return 'pptx', 'ppt_legacy', or 'unknown' based on the first 4 bytes."""
    header = data[:4]
    if header == _MAGIC_PPTX:
        return "pptx"
    if header == _MAGIC_PPT:
        return "ppt_legacy"
    return "unknown"


# ── Public API ─────────────────────────────────────────────────────────────

def extract_ppt_text(file) -> dict:
    """
    Extract all text from a PPTX file.

    Parameters
    ----------
    file : bytes | BinaryIO
        Raw bytes or any file-like object (e.g. Streamlit UploadedFile).

    Returns
    -------
    {
        "slide_count" : int,
        "total_words" : int,
        "slides"      : [{"slide_number": int, "content": str}, ...],
        "full_text"   : str   (all slide text joined — used by the quiz generator)
    }

    Raises
    ------
    UnsupportedFormatError  – not a PPTX (legacy .ppt, PDF, etc.)
    CorruptedFileError      – PPTX header found but file cannot be parsed
    EmptyPresentationError  – opened fine but zero readable text found
    """
    data = file if isinstance(file, bytes) else file.read()

    fmt = _detect_format(data)
    if fmt == "ppt_legacy":
        raise UnsupportedFormatError(
            "Legacy .ppt files are not supported. "
            "Open the file in PowerPoint, save as .pptx, then re-upload."
        )
    if fmt == "unknown":
        raise UnsupportedFormatError(
            "The file does not appear to be a PowerPoint file. "
            "Only .pptx files are accepted."
        )

    # fmt == "pptx" — attempt to open
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise CorruptedFileError(
            "The file could not be opened. "
            f"It may be corrupted or password-protected. ({exc})"
        ) from exc

    if len(prs.slides) == 0:
        raise EmptyPresentationError("The presentation contains no slides.")

    slides = [
        {
            "slide_number": idx,
            "content": _extract_slide_content(slide),
        }
        for idx, slide in enumerate(prs.slides, start=1)
    ]

    # Join only non-empty slides for full_text and word count
    full_text = "\n\n".join(
        s["content"] for s in slides if s["content"].strip()
    )

    if not full_text.strip():
        raise EmptyPresentationError(
            "No readable text was found. "
            "The slides may contain only images or embedded objects."
        )

    return {
        "slide_count": len(slides),
        "total_words": len(full_text.split()),
        "slides": slides,
        "full_text": full_text,
    }


def file_hash(data: bytes) -> str:
    """Return the SHA-256 hex digest of *data* — used as a cache key."""
    return hashlib.sha256(data).hexdigest()


@st.cache_data(show_spinner=False)
def extract_ppt_text_cached(file_bytes: bytes) -> dict:
    """
    Cached version of extract_ppt_text.

    Streamlit caches by the value of *file_bytes*; identical files skip
    re-parsing entirely.  Exceptions propagate to the caller unchanged.
    """
    return extract_ppt_text(file_bytes)


def validate_pptx(file_bytes: bytes) -> tuple[bool, str]:
    """
    Convenience wrapper for the upload screen.

    Returns (True, "") on success or (False, human-readable message) on any error.
    """
    try:
        extract_ppt_text(file_bytes)
        return True, ""
    except PPTParseError as exc:
        return False, str(exc)
    except Exception as exc:
        return False, f"Unexpected error while reading the file: {exc}"


def get_content_preview(slides: list, max_slides: int = 5) -> str:
    """Return a plain-text preview of the first *max_slides* slides."""
    lines = []
    for slide in slides[:max_slides]:
        lines.append(f"── Slide {slide['slide_number']} ──")
        content = slide.get("content", "").strip()
        if content:
            snippet = content[:240]
            if len(content) > 240:
                snippet += "…"
            lines.append(snippet)
        else:
            lines.append("(empty slide — no text)")
        lines.append("")
    return "\n".join(lines).strip()


# ── Private helpers ────────────────────────────────────────────────────────

def _extract_slide_content(slide) -> str:
    """
    Return all text from a slide as a single string.

    Title placeholder text comes first; all other text frames follow in
    document order. Empty shapes are skipped.
    """
    title_parts = []
    body_parts  = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if _is_title_shape(shape):
            title_parts.append(text)
        else:
            body_parts.append(text)

    return "\n".join(title_parts + body_parts)


def _is_title_shape(shape) -> bool:
    """Return True when the shape is a slide title placeholder."""
    try:
        ph = shape.placeholder_format
        if ph is None:
            return False
        return ph.idx == 0 or ph.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        )
    except Exception:
        return False
